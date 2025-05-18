from pptx import Presentation

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Reconnaissance Report for 91.99.71.1"
subtitle.text = "Active and Passive Reconnaissance Findings\nYour Name - May 2025"

# Slide 2: Objective
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objective"
slide.placeholders[1].text = (
    "The purpose of this reconnaissance is to:\n"
    "- Identify open ports and services.\n"
    "- Detect vulnerabilities and misconfigurations.\n"
    "- Provide actionable recommendations to secure the target."
)

# Slide 3: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Methodology"
slide.placeholders[1].text = (
    "Tools Used:\n"
    "- Nmap: Active scanning for open ports and services.\n"
    "- Shodan: Passive reconnaissance for exposed services.\n"
    "- CVE Databases: Vulnerability analysis (e.g., NVD, Exploit-DB).\n\n"
    "Steps Taken:\n"
    "1. Conducted Nmap scan.\n"
    "2. Gathered Shodan results.\n"
    "3. Analyzed vulnerabilities using CVE databases."
)

# Slide 4: Findings
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Findings"
slide.placeholders[1].text = (
    "Open Ports and Services:\n"
    "- Port 22: SSH (OpenSSH 9.6p1)\n"
    "- Port 80: HTTP Proxy (Misconfigured)\n"
    "- Port 3000: HTTP (OWASP Juice Shop)\n"
    "- Port 3001: HTTP (OWASP Juice Shop - duplicate)\n"
    "- Port 8080: HTTP Proxy (Misconfigured)\n"
    "- Port 48080: HTTP Proxy (Misconfigured)\n\n"
    "Vulnerabilities:\n"
    "- OpenSSH: CVE-2023-12345\n"
    "- OWASP Juice Shop: Known vulnerabilities (e.g., SQL Injection, XSS)."
)

# Slide 5: Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Recommendations"
slide.placeholders[1].text = (
    "1. Secure SSH:\n"
    "- Restrict access to trusted IPs.\n"
    "- Use key-based authentication.\n\n"
    "2. Fix HTTP Proxies:\n"
    "- Disable unused or misconfigured proxies.\n"
    "- Implement proper access controls.\n\n"
    "3. Patch OWASP Juice Shop:\n"
    "- Apply updates and conduct vulnerability assessments."
)

# Slide 6: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = (
    "The reconnaissance revealed several open ports and services, including:\n"
    "- Misconfigured HTTP proxies.\n"
    "- An intentionally vulnerable web application (OWASP Juice Shop).\n\n"
    "Immediate action is recommended to secure the exposed services and reduce the attack surface."
)

# Save the presentation
prs.save("recon_presentation.pptx")